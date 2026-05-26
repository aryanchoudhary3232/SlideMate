import io
import re
from typing import List, Optional
from dataclasses import dataclass


@dataclass
class ExtractedPage:
    text: str
    page_number: Optional[int] = None


def _normalize_text(text: str) -> str:
    return _sanitize_text(text)


def _sanitize_text(text: str) -> str:
    # Remove null bytes and control characters
    text = re.sub(r'\x00', '', text)
    text = re.sub(r'[\x01-\x08\x0b\x0c\x0e-\x1f\x7f]', ' ', text)
    text = re.sub(r'\ufffd', ' ', text)
    text = re.sub(r'\s+', ' ', text)
    return text.strip()


def sanitize_text_for_storage(text: str) -> str:
    return _sanitize_text(text)


def extract_text_from_buffer(
    buffer: bytes,
    mime_type: str,
    file_name: str,
) -> List[ExtractedPage]:
    fname = file_name.lower()

    if "pdf" in mime_type or fname.endswith(".pdf"):
        return _extract_pdf(buffer)

    if "wordprocessingml" in mime_type or fname.endswith(".docx"):
        return _extract_docx(buffer)

    if "presentationml" in mime_type or fname.endswith(".pptx"):
        return _extract_pptx(buffer)

    if mime_type.startswith("text/") or fname.endswith(".txt"):
        return [ExtractedPage(text=_normalize_text(buffer.decode("utf-8", errors="replace")))]

    raise ValueError("Unsupported file type. Upload PDF, PPTX, DOCX, or TXT.")


def _extract_pdf(buffer: bytes) -> List[ExtractedPage]:
    import PyPDF2
    reader = PyPDF2.PdfReader(io.BytesIO(buffer))
    pages = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        pages.append(ExtractedPage(text=_normalize_text(text), page_number=i))
    return pages


def _extract_docx(buffer: bytes) -> List[ExtractedPage]:
    import docx
    doc = docx.Document(io.BytesIO(buffer))
    text = "\n".join(p.text for p in doc.paragraphs)
    return [ExtractedPage(text=_normalize_text(text))]


def _extract_pptx(buffer: bytes) -> List[ExtractedPage]:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(buffer))
    pages = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs)
                    if line.strip():
                        texts.append(line)
        pages.append(ExtractedPage(
            text=_normalize_text(" ".join(texts)),
            page_number=i,
        ))
    return pages


def chunk_extracted_pages(pages: List[ExtractedPage]):
    """Split pages into overlapping word-level chunks."""
    chunks = []
    chunk_index = 0
    size = 650
    overlap = 90

    for page in pages:
        words = [w for w in page.text.split() if w]
        start = 0
        while start < len(words):
            content = sanitize_text_for_storage(" ".join(words[start: start + size]))
            if len(content) >= 40:
                chunks.append({
                    "content": content,
                    "page_number": page.page_number,
                    "chunk_index": chunk_index,
                })
                chunk_index += 1
            start += size - overlap

    return chunks
